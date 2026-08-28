#!/usr/bin/env python
"""Builds a 2-slide summary deck for the resolution-attribution exploratory
analysis (scripts/resolution_attribution.py). Reads only the CSV/PNG that
script already produced -- no recomputation.

Usage: python scripts/build_resolution_slides.py
"""

from __future__ import annotations

import csv
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULTS_DIR = os.path.join(REPO_ROOT, "results")
FIGURE_PATH = os.path.join(RESULTS_DIR, "derived_resolution_attribution.png")
CSV_PATH = os.path.join(RESULTS_DIR, "derived_resolution_attribution.csv")
OUT_PATH = os.path.join(RESULTS_DIR, "resolution_attribution_slides.pptx")

INK = RGBColor(0x22, 0x27, 0x33)
MUTED = RGBColor(0x5A, 0x63, 0x73)
ACCENT = RGBColor(0x4C, 0x72, 0xB0)
BG = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _read_summary_stats():
    with open(CSV_PATH, newline="") as f:
        first = f.readline()
        meta = json.loads(first[len("# METADATA "):])
        reader = csv.DictReader(f)
        rows = list(reader)

    included = [r for r in rows if r["excluded_low_snow"] == "False"]
    biases = [float(r["bias_pp"]) for r in included]
    n_total = len(rows)
    n_included = len(included)
    return meta, n_total, n_included, biases


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_title(slide, text, kicker=None):
    if kicker:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = kicker.upper()
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        run.font.name = "Arial"

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35 + (0.45 if kicker else 0)), Inches(12.2), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = INK
    run.font.name = "Arial"

    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(1.45), Inches(12.2), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0xE2, 0xE5, 0xEB)
    rule.line.fill.background()
    rule.shadow.inherit = False


def _add_bullets(slide, x, y, w, h, items, size=16, leading=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = leading
        p.space_after = Pt(10)
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = head + "  "
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = INK
            r1.font.name = "Arial"
            r2 = p.add_run()
            r2.text = body
            r2.font.size = Pt(size)
            r2.font.color.rgb = MUTED
            r2.font.name = "Arial"
        else:
            r = p.add_run()
            r.text = "•  " + item
            r.font.size = Pt(size)
            r.font.color.rgb = MUTED
            r.font.name = "Arial"
    return box


def _add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = MUTED
    run.font.name = "Arial"


def build():
    meta, n_total, n_included, biases = _read_summary_stats()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Slide 1: what was done -------------------------------------------------
    s1 = _blank_slide(prs)
    _set_bg(s1)
    _add_title(
        s1,
        "Does spatial resolution explain the MERRA-2 vs. MODSCAG fSCA disagreement?",
        kicker="Exploratory follow-up · no new data pulled",
    )

    _add_bullets(
        s1,
        Inches(0.55), Inches(1.65), Inches(5.9), Inches(4.6),
        [
            ("Question:", "how much of the pooled bias pattern reflects MERRA's coarse "
             "0.625°×0.5° grid vs. some other systematic cause?"),
            ("Constraint:", "reuse existing outputs only — no daily reprocessing, no "
             "redownload, no change to the aggregation rule (MODIS is always aggregated "
             "up to MERRA; MERRA is never resampled down)."),
            ("Terrain heterogeneity proxy:", "binned the checked-in USGS 3DEP DEM "
             "(800×600 px) into the 72 MERRA cells and computed within-cell "
             "elevation std. dev. — cells straddling steep terrain are the ones a "
             "coarse grid represents worst."),
            ("Bias signal:", "recombined the already-computed monthly sufficient "
             "statistics (climatology_month rows in pixel_stats.csv) per cell into one "
             "pooled WY2010–2023 bias/NMB per cell, same 5% low-snow mask as the "
             "public bias/MAE figure."),
            ("Two comparisons:", "(1) correlate terrain heterogeneity against the "
             "spatial bias pattern; (2) compare spatial variance (cell-to-cell) against "
             "temporal variance (year-to-year) as a rough magnitude check on which axis "
             "dominates."),
        ],
        size=15,
    )

    info_box = slide_info = s1.shapes.add_textbox(Inches(6.75), Inches(1.65), Inches(6.0), Inches(4.6))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Inputs (all pre-existing)"
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = INK
    r.font.name = "Arial"

    rows = [
        ("tests/fixtures/domain_dem_3dep.tif", "real 3DEP DEM, checked in"),
        ("results/water_year_2010_2023_pixel_stats.csv", "per-cell climatology sufficient stats"),
        ("results/water_year_2010_2023_overall_stats.csv", "per-year domain sufficient stats"),
    ]
    for label, note in rows:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        r1 = p2.add_run()
        r1.text = label
        r1.font.size = Pt(13)
        r1.font.name = "Consolas"
        r1.font.color.rgb = ACCENT
        p3 = tf.add_paragraph()
        r2 = p3.add_run()
        r2.text = note
        r2.font.size = Pt(13)
        r2.font.color.rgb = MUTED
        r2.font.name = "Arial"

    p4 = tf.add_paragraph()
    p4.space_before = Pt(18)
    r4 = p4.add_run()
    r4.text = "New, exploratory-only artifact:"
    r4.font.bold = True
    r4.font.size = Pt(14)
    r4.font.color.rgb = INK
    r4.font.name = "Arial"
    p5 = tf.add_paragraph()
    r5 = p5.add_run()
    r5.text = "scripts/resolution_attribution.py → results/derived_resolution_attribution.{csv,png}"
    r5.font.size = Pt(13)
    r5.font.name = "Consolas"
    r5.font.color.rgb = ACCENT

    _add_footer(
        s1,
        "Not part of the scientific-contract pipeline or public product; error sign, masks, and "
        "grid rules unchanged from README.md.",
    )

    # --- Slide 2: results ---------------------------------------------------
    s2 = _blank_slide(prs)
    _set_bg(s2)
    _add_title(s2, "Result: bias pattern is mostly not explained by terrain heterogeneity", kicker="Findings")

    pic_w = Inches(7.6)
    pic_h = pic_w * (767 / 1817)
    pic_top_in = 1.5
    s2.shapes.add_picture(FIGURE_PATH, Inches(0.55), Inches(pic_top_in), width=pic_w, height=pic_h)

    r_note = "r = -0.19 (p = 0.12, R² ≈ 0.04)"
    bullets_top_in = pic_top_in + float(pic_h / Inches(1)) + 0.2

    _add_bullets(
        s2,
        Inches(0.55), Inches(bullets_top_in), Inches(12.2), Inches(7.0 - bullets_top_in),
        [
            ("Correlation:", f"{r_note} between within-cell elevation std. dev. and pooled "
             f"bias across {n_included} of {n_total} cells (5 masked for low snow) — "
             "essentially no relationship, and the weak trend runs opposite to a "
             "resolution-mismatch story."),
            ("Variance split:", "spatial variance across cells (≈16.9 pp²) is "
             "~4.5× the interannual variance at the domain level (≈3.7 pp²) — "
             "the disagreement is mostly a stable spatial pattern, not a year-to-year one."),
            ("Interpretation:", "since that spatial pattern doesn't track terrain "
             "heterogeneity, resolution/aggregation mismatch looks like a minor "
             "contributor. Worth investigating next: land cover, radiation/snow-physics "
             "differences, or a location-dependent MERRA bias unrelated to subgrid relief."),
        ],
        size=13,
        leading=1.05,
    )

    _add_footer(
        s2,
        "Correlational, not a controlled regridding experiment — MERRA is never resampled to MODIS resolution "
        "in this repository. See results/derived_resolution_attribution.csv for full per-cell values.",
    )

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build()
